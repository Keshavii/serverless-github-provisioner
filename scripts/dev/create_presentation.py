#!/usr/bin/env python3
"""
GitHub Repository Automation - PowerPoint Generator
Run: python3 create_presentation.py
Output: GitHub_Automation_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_bullets(tf, items, size=18):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0, 82, 204)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(128, 128, 128)
GREEN = RGBColor(0, 176, 80)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
box.text_frame.text = "GitHub Repository Automation"
box.text_frame.paragraphs[0].font.size = Pt(48)
box.text_frame.paragraphs[0].font.bold = True
box.text_frame.paragraphs[0].font.color.rgb = BLUE
box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
box.text_frame.text = "From 30 Minutes to 90 Seconds"
box.text_frame.paragraphs[0].font.size = Pt(28)
box.text_frame.paragraphs[0].font.color.rgb = DARK
box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 2: Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Problem: Manual Repository Creation"
tf = slide.placeholders[1].text_frame
tf.text = "Current Process: 13 Manual Steps"
add_bullets(tf, ["TIME: 30-45 min/repo", "CAPACITY: 20 repos/day", "ERROR RATE: 10-20%", "COST: $36,000/year", "BOTTLENECK: 2-week backlog"], 20)

# Slide 3: Pain
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Pain Quantified"
tf = slide.placeholders[1].text_frame
tf.text = "Annual Impact:"
add_bullets(tf, ["480 hours/year = 3 MONTHS", "Labor: $36,000/year", "Errors: $5,000/year", "Total: $41,000/year burned"], 18)

# Slide 4: Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Solution: Serverless Automation"
tf = slide.placeholders[1].text_frame
tf.text = "AUTOMATED: 1 STEP"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN
p = tf.add_paragraph()
p.text = "1. Developer creates JIRA ticket"
p.font.size = Pt(18)
p = tf.add_paragraph()
p.text = "↓ 90 seconds ↓"
p.font.italic = True
p.alignment = PP_ALIGN.CENTER
add_bullets(tf, ["Repo created", "JIRA updated", "Ticket closed", "Dev notified"], 18)

# Slide 5: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Event-Driven Serverless Architecture"
tf = slide.placeholders[1].text_frame
tf.text = "Flow:"
add_bullets(tf, ["JIRA → Lambda → SQS → Lambda → GitHub", "Failure → DLQ → Alert", "Supporting: Secrets Manager, CloudWatch"], 16)

# Slide 6: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Technology Stack"
tf = slide.placeholders[1].text_frame
tf.text = "AWS: Lambda, SQS, Secrets Manager"
add_bullets(tf, ["Python 3.11", "PyGithub 2.1.1", "Pydantic validation", "Terraform IaC"], 18)

# Slide 7: Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Production Features"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["3 automatic retries", "Idempotent operations", "GitHub App auth", "Encrypted secrets", "Full audit trail"], 20)

# Slide 8: Time Savings
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Time Savings"
tf = slide.placeholders[1].text_frame
tf.text = "95% REDUCTION"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
add_bullets(tf, ["Manual: 30 min", "Automated: 90 sec", "", "480 hours/year saved", "= 60 work days"], 22)

# Slide 9: Cost & ROI
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Cost Savings & ROI"
tf = slide.placeholders[1].text_frame
tf.text = "MANUAL: $41,000/year"
add_bullets(tf, ["AUTOMATED: $29/year", "Net savings: $40,971/year", "ROI: 141,000%"], 20)

# Slide 10: Quality
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Zero Human Errors"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["Wrong visibility: 10% → 0%", "Missing metadata: 20% → 0%", "Typos: 10% → 0%", "", "100% IMPROVEMENT"], 20)

# Slide 11: Scalability
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Built to Scale 10x"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["TODAY: 60 repos → $2/month", "6 MONTHS: 200 repos → $2/month", "1 YEAR: 500 repos → $3/month", "10x: 6,000 repos → $10/month", "", "Manual at 500: Need 5 engineers"], 18)

# Slide 12: Comparison
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Manual vs Automated"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["Time: 30 min vs 90 sec ✅", "Capacity: 60 vs Unlimited ✅", "Errors: 10-20% vs 0% ✅", "Cost: $3,000 vs $2 ✅", "", "WINNER: AUTOMATION"], 18)

# Slide 13: Production Ready
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Production Ready"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["✅ Code quality", "✅ Infrastructure (Terraform)", "✅ Security (encrypted)", "✅ Observability (CloudWatch)", "✅ 25+ test repos created"], 20)

# Slide 14: Deployment
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "4-Week Rollout Plan"
tf = slide.placeholders[1].text_frame
add_bullets(tf, ["WEEK 1: Test ✅", "WEEK 2: Parallel testing", "WEEK 3: Production deployment", "WEEK 4: Stabilization"], 18)

# Slide 15: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Recommendation"
tf = slide.placeholders[1].text_frame
tf.text = "APPROVE PRODUCTION DEPLOYMENT"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN
add_bullets(tf, ["", "$41,000/year savings", "480 hours/year freed", "0% error rate", "Production ready NOW"], 20)

# Save
filename = "GitHub_Automation_Presentation.pptx"
prs.save(filename)
print(f"✅ Created: {filename}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"\nTo view: open {filename}")

